# -*- coding: utf-8 -*-
"""生成《肠道花园 Gut Garden》核心代码 PPT（10 页 16:9，每页全是代码）。
代码从源码按行提取，去掉头部注释，每页放一个模块的核心函数，填满页面。
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
GREEN = RGBColor(0x2E, 0x5B, 0x2A)
DGREEN = RGBColor(0x1F, 0x3D, 0x2E)
GRAY = RGBColor(0x88, 0x88, 0x88)
CODE_BG = RGBColor(0xF6, 0xF8, 0xF2)
CODE_LN = RGBColor(0xB7, 0xC4, 0xAE)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def source_lines(rel):
    with open(os.path.join(ROOT, rel), 'r', encoding='utf-8') as f:
        return f.read().split('\n')


def extract(rel, start, end):
    """取第 start..end 行，去掉头部注释/空行，返回行列表。"""
    lines = source_lines(rel)[start - 1:end]
    i = 0
    while i < len(lines):
        s = lines[i].strip()
        if s == '' or s.startswith('/**') or s.startswith('/*') or s.startswith('*') or s.startswith('*/'):
            i += 1
        else:
            break
    lines = lines[i:]
    while lines and not lines[-1].strip():
        lines = lines[:-1]
    return lines


def join_blocks(*blocks):
    out = []
    for b in blocks:
        if out:
            out.append('')
        out.extend(b)
    return out


def style_run(run, size, bold=False, color=DGREEN, latin='Consolas', ea='微软雅黑'):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn('a:ea'))
    if ea_el is None:
        ea_el = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea_el)
    ea_el.set('typeface', ea)


def code_slide(num, title, subtitle, lines):
    slide = prs.slides.add_slide(BLANK)
    # 标题
    add = slide.shapes.add_textbox(Inches(0.5), Inches(0.32), Inches(12.3), Inches(0.55))
    tf = add.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = f'{num}  {title}'
    style_run(r, 20, bold=True, color=GREEN, latin='Calibri')
    # 副标题（文件定位）
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.35))
    tf = sub.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = subtitle
    style_run(r, 11, color=GRAY, latin='Consolas')
    # 代码框
    n = len(lines)
    size = max(9, min(13, int(405 / (1.18 * n))))
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(1.22), Inches(12.63), Inches(5.95))
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.color.rgb = CODE_LN
    box.line.width = Pt(1)
    box.shadow.inherit = False
    box.adjustments[0] = 0.02
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)
    first = True
    for line in lines:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.line_spacing = 1.0
        run = para.add_run()
        run.text = line if line else ' '
        style_run(run, size)
    # 页码
    pg = slide.shapes.add_textbox(Inches(12.4), Inches(7.18), Inches(0.8), Inches(0.3))
    tf = pg.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    r = tf.paragraphs[0].add_run()
    r.text = f'{num} / 10'
    style_run(r, 9, color=GRAY, latin='Consolas')
    return slide


# ============ 1. AI 知识底座注入 ============
code_slide(1, 'AI 导游 · 知识底座注入',
           'server/src/modules/ai/knowledge-base.ts — buildKnowledgeText()',
           extract('server/src/modules/ai/knowledge-base.ts', 6, 43))

# ============ 2. AI 流式输出与三级降级 ============
code_slide(2, 'AI 导游 · 流式输出与三级降级',
           'server/src/modules/ai/ai.routes.ts buildMessages() + ai-stream.ts streamAiChunks()',
           join_blocks(
               extract('server/src/modules/ai/ai.routes.ts', 11, 23),
               extract('server/src/modules/ai/ai-stream.ts', 50, 69),
           ))

# ============ 3. 便便垂直 AI · 角色约束 ============
code_slide(3, '便便垂直 AI · 角色与输出约束',
           'server/src/modules/stool/stool-ai.ts — STOOL_AI_SYSTEM_PROMPT',
           extract('server/src/modules/stool/stool-ai.ts', 13, 45))

# ============ 4. 便便垂直 AI · 生成建议 ============
code_slide(4, '便便垂直 AI · 生成建议（含降级）',
           'server/src/modules/stool/stool-ai.ts — generateStoolSuggestion()',
           extract('server/src/modules/stool/stool-ai.ts', 80, 115))

# ============ 5. 每日打卡 · 任务定义与确认 ============
code_slide(5, '每日打卡 · 任务定义与确认',
           'server/src/modules/checkin/checkin.service.ts — TASK_DEFS + confirmTask()',
           join_blocks(
               extract('server/src/modules/checkin/checkin.service.ts', 9, 19),
               extract('server/src/modules/checkin/checkin.service.ts', 178, 196),
           ))

# ============ 6. 每日打卡 · 连续打卡计算 ============
code_slide(6, '每日打卡 · 连续打卡计算',
           'server/src/modules/checkin/checkin.service.ts — computeStreaks()',
           extract('server/src/modules/checkin/checkin.service.ts', 60, 96))

# ============ 7. 徽章引擎 · 规则即数据 ============
code_slide(7, '徽章引擎 · 规则即数据',
           'server/src/modules/badges/engine.ts — checkRule()',
           extract('server/src/modules/badges/engine.ts', 154, 190))

# ============ 8. 花园成长 · 阶段规则 ============
code_slide(8, '花园成长 · 阶段规则数据化',
           'server/src/modules/garden/stage.service.ts — STAGE_REQS + evaluateStage()',
           join_blocks(
               extract('server/src/modules/garden/stage.service.ts', 17, 24),
               extract('server/src/modules/garden/stage.service.ts', 47, 58),
           ))

# ============ 9. 花园投喂 · 前端交互 ============
code_slide(9, '花园投喂 · 前端拖拽反馈',
           'web/src/hooks/useFeedLogic.ts — FOOD_EFFECTS + handleDrop()（注册用户分支）',
           join_blocks(
               extract('web/src/hooks/useFeedLogic.ts', 8, 16),
               extract('web/src/hooks/useFeedLogic.ts', 29, 53),
           ))

# ============ 10. 验证码登录 ============
code_slide(10, '验证码登录 · 模拟短信',
           'server/src/modules/auth/codeStore.ts + auth.service.ts loginWithCode()',
           join_blocks(
               extract('server/src/modules/auth/codeStore.ts', 1, 24),
               extract('server/src/modules/auth/auth.service.ts', 10, 24),
           ))

OUT = r'D:\GutGardenBeta\docs\比赛方案_核心代码整理.pptx'
prs.save(OUT)
print('已生成:', OUT, '共', len(prs.slides.__iter__.__self__._sldIdLst), '页')
